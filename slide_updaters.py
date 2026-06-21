"""Her slayt için güncelleme fonksiyonları."""
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from slide_utils import set_para_text, set_bullet_para
from data_sources import fmt_pct, fmt_mtl

NS_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

def _pct_str(v, decimals=1):
    """0.602 → '60,2'"""
    if v is None: return "?"
    return f"{v*100:.{decimals}f}".replace('.', ',')

def _pct_str_sign(v, decimals=1):
    if v is None: return "?"
    s = f"{abs(v)*100:.{decimals}f}".replace('.', ',')
    return f"+{s}" if v >= 0 else f"-{s}"

def _set_run_text(para, text, bold=None, color=None):
    """Paragrafı tamamen temizleyip tek run yazar."""
    _clear_para(para)
    if not text: return
    run = para.add_run()
    run.text = text
    if bold is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color

def _left(sh): return Emu(sh.left).inches if sh.left else 0
def _w(sh):    return Emu(sh.width).inches if sh.width else 0

def _clear_para(para):
    """Paragraftaki tüm run'ları siler."""
    for r_el in list(para._p.findall(f'{NS_A}r')):
        para._p.remove(r_el)

def _set_run_text(para, text, bold=None, color=None, size_pt=None):
    """Paragrafı tamamen temizleyip tek run yazar."""
    _clear_para(para)
    if not text: return
    run = para.add_run()
    run.text = text
    if bold is not None:   run.font.bold = bold
    if color is not None:  run.font.color.rgb = color
    if size_pt is not None: run.font.size = Pt(size_pt)

def _update_cat_run(run, new_text):
    """Run metnini günceller, formatı korur."""
    run.text = new_text


# ── Tarih metinleri güncelleme ───────────────────────────────────────────────
def replace_date_texts(slides, old_texts, new_texts):
    """Tüm slaytlarda eski tarih/ay yazılarını yenisiyle değiştirir.
    Unicode apostrof varyantlarını da dikkate alır."""
    pairs = list(zip(old_texts, new_texts))
    if not pairs: return
    # Her eski metni hem ASCII (') hem Unicode (\u2019) apostrof ile dene
    expanded = []
    for old, new in pairs:
        expanded.append((old, new))
        if "'" in old:
            expanded.append((old.replace("'", "\u2019"), new))
        if "\u2019" in old:
            expanded.append((old.replace("\u2019", "'"), new))

    for slide in slides:
        for sh in slide.shapes:
            if not sh.has_text_frame: continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in expanded:
                        if old and old in run.text:
                            run.text = run.text.replace(old, new)


# ── Slayt 2: Executive Summary ───────────────────────────────────────────────
def update_slide2(slide, rapor, cevirdigim, kpi_manual, config):
    """Slayt 2 — tüm değerleri Haziran/rapor verileriyle günceller.
    Eski içerikten kalan run'ları temizler."""
    ay_sezon = cevirdigim.get('ay_sezon_hedef', {})
    dr       = config.get('date_range', '')
    old_dr   = config.get('old_date_range', '')
    old_ay   = config.get('old_ay_etiket', '')
    new_ay   = config.get('month_label', '')
    old_month = old_ay.split("'")[0] if old_ay else ''
    new_month = new_ay.split("'")[0] if new_ay else ''

    # Paragraftaki sayı run'larını temizleyip ilk run'a yeni değeri yazar.
    # Etiket içeren run'ları (harf-baskın) korur.
    def _set_value_keep_labels(para, new_value):
        runs = para.runs
        if not runs: return
        wrote = False
        for r in runs:
            txt = r.text
            # Bu run sayı/birim mi yoksa etiket mi?
            has_digit = any(c.isdigit() for c in txt)
            has_letter = any(c.isalpha() for c in txt)
            # "MTL" / " MTL" / " TL" gibi birim suffix'leri sayı sayılır
            is_unit_suffix = txt.strip() in ('MTL', 'TL', '%', '')
            # Etiket = harf var + sayı yok (MTD Hedef, MTD Gerçekleşen, HGO gibi)
            is_label = has_letter and not has_digit and not is_unit_suffix
            
            if is_label:
                continue  # etiketi koru
            # Sayı/birim run'ı
            if not wrote:
                r.text = new_value
                r.font.bold = True
                wrote = True
            else:
                r.text = ''
        if not wrote:
            # Hiç sayı run'ı yoktu - ilk run'a ekle
            runs[0].text = new_value + runs[0].text

    # %-bölünmüş paragraf güncelleme ('%' + '26,8' yapısı)
    def _set_pct(para, new_val):
        new_val = str(new_val).strip()
        runs = para.runs
        if not runs: return
        # Tüm sayısal run'ları topla, ilkini güncelle, diğerlerini temizle
        _set_value_keep_labels(para, new_val)

    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text.strip()
        l = _left(sh)
        paras = sh.text_frame.paragraphs

        # Sol büyük kutu: MTD Hedef + MTL + MTD Gerçekleşen + HGO
        if 'MTD Hedef' in t and 'MTL' in t and l < 4:
            hgo = rapor.get('hgo')
            hgo_str = f"{hgo:.1f}%".replace('.', ',') if hgo else '-'
            # Paragrafları işle
            for pi, para in enumerate(paras):
                pt = ''.join(r.text for r in para.runs)
                if not pt.strip(): continue
                # Hedef paragrafı: "MTL" var, etiket yok ya da ilk paragraf
                if 'MTL' in pt and ('Hedef' not in pt and 'Gerçekleşen' not in pt):
                    _set_value_keep_labels(para, fmt_mtl(rapor['hedef_mtl']))
                elif 'MTL' in pt and 'Gerçekleşen' in pt:
                    # P3: sayı + 'MTD Gerçekleşen' etiketi
                    _set_value_keep_labels(para, fmt_mtl(rapor['grc_mtl']))
                elif '%' in pt and 'HGO' not in pt and 'MTL' not in pt:
                    # HGO sayısı (109,4% gibi)
                    _set_value_keep_labels(para, hgo_str)
                elif 'Hedef' in pt and 'MTL' not in pt:
                    # Sadece "MTD Hedef" etiketi - dokunma
                    pass

        # 6 küçük KPI kutusu
        elif 'Net TL' in t and l > 4:
            v = rapor.get('net_tl_vs_ly', '')
            if v: _set_value_keep_labels(paras[0], str(v))
        elif 'Net Adet' in t and l > 4:
            v = rapor.get('net_adet_vs_ly', '')
            if v: _set_value_keep_labels(paras[0], str(v))
        elif 'İşlem Adet' in t:
            v = rapor.get('islem_adet_vs_ly', '')
            if v: _set_value_keep_labels(paras[0], str(v))
        elif 'UPT' in t and l > 7:
            v = rapor.get('upt_vs_ly') or kpi_manual.get('upt', '')
            if v: _set_value_keep_labels(paras[0], str(v))
        elif 'Sepet' in t and 'TL' in t:
            v = rapor.get('sepet_tl_ty') or kpi_manual.get('sepet_tl', '')
            if v:
                # Run yapısını koru: run0=sayı(bold), run1=' TL'(bold=False)
                runs = paras[0].runs
                if len(runs) >= 2:
                    runs[0].text = str(v)
                    # runs[1] (' TL' suffix) dokunulmadan kalır - bold/font korunur
                else:
                    _set_value_keep_labels(paras[0], f"{v} TL")
        elif 'Birim Fiyat' in t and l > 5:
            v = rapor.get('birim_fiyat_ty') or kpi_manual.get('birim_fiyat', '')
            if v:
                runs = paras[0].runs
                if len(runs) >= 2:
                    runs[0].text = str(v)
                else:
                    _set_value_keep_labels(paras[0], f"{v} TL")

        # Ay başlığı
        elif 'MTD Hedef' in t and 'Performans' in t:
            for para in paras:
                for run in para.runs:
                    if old_month and old_month in run.text:
                        run.text = run.text.replace(old_month, new_month)
                    if old_dr and old_dr in run.text:
                        run.text = run.text.replace(old_dr, dr)

        # Denim kutusu — KATEGORİ var + NON yok = Denim Kategori kutusu
        elif 'KATEGOR' in t.upper() and 'NON' not in t.upper() and ('DENIM' in t.upper() or 'DENİM' in t.upper()):
            k_val = ay_sezon.get('WomenDenim All')
            e_val = ay_sezon.get('MenDenim All')
            for pi, para in enumerate(paras):
                runs = para.runs
                if not runs: continue
                full_p = ''.join(r.text for r in runs)
                # Denim All içeren satırlar
                if 'Denim' in full_p and 'All' in full_p:
                    # Kadın mı erkek mi? Önceki paragraflara bak
                    is_kadin = False
                    for back in range(pi-1, -1, -1):
                        bt = paras[back].text.strip().upper()
                        if 'KADIN' in bt: is_kadin = True; break
                        if 'ERKEK' in bt: is_kadin = False; break
                    val = k_val if is_kadin else e_val
                    if val is not None and runs:
                        from pptx.dml.color import RGBColor
                        runs[-1].text = _pct_str(val)
                        runs[-1].font.bold = True
                        runs[-1].font.color.rgb = RGBColor(0,128,0) if val >= 0 else RGBColor(192,0,0)

        # NonDenim kutuları
        elif 'NON' in t.upper() and 'KATEGOR' in t.upper():
            if 'ERKEK' in t.upper():
                cats = {'Aksesuar': ay_sezon.get('MenAccessories'),
                        'Non Denim Alt': ay_sezon.get('MenNon Denim Bottoms'),
                        'Non': ay_sezon.get('MenNon Denim Bottoms'),
                        'Ceket': ay_sezon.get('MenJackets'),
                        'Gömlek': ay_sezon.get('MenShirts'),
                        'Tees': ay_sezon.get('MenTees')}
            else:
                cats = {'Aksesuar': ay_sezon.get('WomenAccessories'),
                        'Non Denim Alt': ay_sezon.get('WomenNon Denim Bottoms'),
                        'Non': ay_sezon.get('WomenNon Denim Bottoms'),
                        'Ceket': ay_sezon.get('WomenJackets'),
                        'Gömlek': ay_sezon.get('WomenShirts'),
                        'Tees': ay_sezon.get('WomenTees')}
            _update_nondenim_box_inplace(sh, cats)


def _update_nondenim_box_inplace(sh, cats):
    """NonDenim kutusu - paragraf yapisini koruyarak SON run'daki sayiyi gunceller.
    Isaret ayri run'daysa onu da temizler (Tees gibi cok run'lu durumlar icin)."""
    paras = sh.text_frame.paragraphs
    for para in paras:
        runs = para.runs
        if not runs: continue
        full = ''.join(r.text for r in runs).strip()
        if not full: continue
        if 'Hedefin' in full or 'KATEGORI' in full.upper() or 'KATEGORİ' in full.upper() or full.upper() in ('KADIN','ERKEK'):
            continue
        cat_text = ''.join(r.text for r in runs[:-1]).strip()
        matched = None
        for key, val in cats.items():
            if key.lower() in cat_text.lower():
                matched = val; break
        if matched is None: continue
        new_val = _pct_str_sign(matched)
        from pptx.dml.color import RGBColor
        color = RGBColor(0,128,0) if matched >= 0 else RGBColor(192,0,0)
        # Sadece tam olarak '+' veya '-' olan run'ları temizle (boşluk run'larına dokunma)
        for ri in range(len(runs)-2, -1, -1):
            rt = runs[ri].text.strip()
            if rt in ('+', '-'):
                runs[ri].text = ''
            elif rt == '':
                continue  # boşluk run'u - koru, atlama devam et
            else:
                break
        runs[-1].text = new_val
        runs[-1].font.bold = True
        runs[-1].font.color.rgb = color



# Kategori → Furkan anahtar haritaları
_ERKEK_CAT_MAP = {
    'Denim':        'MenDenim All',
    'Ceket':        'MenJackets',
    'Gömlek':       'MenShirts',
    'Non':          'MenNon Denim Bottoms',
    'Aksesuar':     'MenAccessories',
    'Tees':         'MenTees',
}
_KADIN_CAT_MAP = {
    'Denim':        'WomenDenim All',
    'Ceket':        'WomenJackets',
    'Gömlek':       'WomenShirts',
    'Non':          'WomenNon Denim Bottoms',
    'Aksesuar':     'WomenAccessories',
    'Tees':         'WomenTees',
}

def update_slide6(slide, cevirdigim, ocr_data, config):
    _update_lfl_cat_box(slide, _ERKEK_CAT_MAP, cevirdigim)

def update_slide8(slide, cevirdigim, ocr_data, config):
    _update_lfl_cat_box(slide, _KADIN_CAT_MAP, cevirdigim)


# ── Slayt 3: Hedef Değerlendirme ─────────────────────────────────────────────
def update_slide3(slide, rapor, cevirdigim, kpi_manual, config):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text; l = _left(sh)
        if l >= 6 and 'Hedef:' in t and 'Gerçekleşme' in t:
            hgo = rapor.get('hgo')
            hgo_str = f'%{hgo:.1f}'.replace('.', ',') if hgo else '-'
            set_bullet_para(sh, 0, 'Hedef: ', fmt_mtl(rapor['hedef_mtl']), value_bold=True)
            set_bullet_para(sh, 1, 'Hedef Gerçekleşme:  ', fmt_mtl(rapor['grc_mtl']), value_bold=True)
            set_bullet_para(sh, 2, 'Hedef Gerçekleşme Oranı: ', hgo_str, value_bold=True)
        elif l >= 6 and 'LFL' in t:
            na = rapor.get('net_adet_vs_ly', '')
            ia = rapor.get('islem_adet_vs_ly', '')
            bf_ty = rapor.get('birim_fiyat_ty') or kpi_manual.get('birim_fiyat', '')
            bf_ly = rapor.get('birim_fiyat_vs_ly', '')
            uv = rapor.get('upt_vs_ly') or kpi_manual.get('upt', '')
            ut = rapor.get('upt_ty', '')
            p0 = f"LFL net adette {na} işlem adette ise {ia} büyüme gerçekleşmiştir."
            if bf_ly and bf_ty:
                p1 = f"Birim fiyat artışı {bf_ly} gerçekleşerek {bf_ty} TL'ye gelmiştir."
            elif bf_ty:
                p1 = f"Birim fiyat {bf_ty} TL'ye gelmiştir."
            else: p1 = ''
            p2 = f"UPT geçen seneye göre {uv} büyüme ile {ut} seviyesine yerleşmiştir." if uv else ''
            set_para_text(sh, 0, p0)
            if p1: set_para_text(sh, 1, p1)
            if p2: set_para_text(sh, 2, p2)


# ── Slayt 4: Yönetici Özeti ──────────────────────────────────────────────────
def update_slide4(slide, rapor, cevirdigim, config):
    dr  = config.get('date_range', '')
    ay  = rapor.get('ay_etiket', '')
    hgo = rapor.get('hgo')
    hgo_str = f'%{hgo:.1f}'.replace('.', ',') if hgo else '-'
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text
        if 'MTL' in t and len(t) > 80:
            new = (f"Genel {ay} ayında {fmt_mtl(rapor['hedef_mtl'])} olan hedefin {dr} tarihleri "
                   f"arasında {fmt_mtl(rapor['grc_mtl'])} olarak gerçekleşmesiyle birlikte "
                   f"HGO {hgo_str} olarak gerçekleşmiştir.")
            set_para_text(sh, 0, new)
            paras = sh.text_frame.paragraphs
            for pi in range(1, len(paras)):
                ptext = paras[pi].text.strip()
                if not ptext: continue
                if 'Best Seller' in ptext or 'jean' in ptext.lower():
                    break
                if any(kw in ptext for kw in ['Mayıs', 'Haziran', 'Kurban', 'tatil', 'bayram',
                                                'kategorisinde', 'kategorisinin', 'kategori',
                                                'Ürüne', 'indirim', 'İşlem Adette', 'hedefimiz',
                                                "LFL'de", "Store'da", 'oranında']):
                    _clear_para(paras[pi])
            break


# ── Slayt 5: LFL Toplam ──────────────────────────────────────────────────────
def update_slide5(slide, cevirdigim, config):
    grc = cevirdigim.get('gerceklesen_mtd', {})
    hdf = cevirdigim.get('hedef_mtd', {})
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text
        if ('Sezon+Devam' in t or ('Hedef:' in t and 'Gerçekleşen:' in t)) and _left(sh) > 6:
            gt_g = _pct_str(grc.get('Genel Toplam'))
            gt_h = _pct_str(hdf.get('Genel Toplam')) if hdf.get('Genel Toplam') else '?'
            er_g = _pct_str(grc.get('Toplam Men'))
            er_h = _pct_str(hdf.get('Toplam Men')) if hdf.get('Toplam Men') else '?'
            kd_g = _pct_str(grc.get('Toplam Women'))
            kd_h = _pct_str(hdf.get('Toplam Women')) if hdf.get('Toplam Women') else '?'
            new = (f"Sezon+Devam+Gelecek Toplam\nHedef: {gt_h}\nGerçekleşen: {gt_g}\n"
                   f"Erkek Kategorisi \nHedef: {er_h}\nGerçekleşen: {er_g}\n"
                   f"Kadın Kategorisi \nHedef: {kd_h}\nGerçekleşen: {kd_g}")
            _set_run_text(sh.text_frame.paragraphs[0], new)
            for pi in range(1, len(sh.text_frame.paragraphs)):
                _clear_para(sh.text_frame.paragraphs[pi])
            break


def _update_lfl_cat_box(slide, cat_map, cevirdigim):
    grc = cevirdigim.get('gerceklesen_mtd', {})
    hdf = cevirdigim.get('hedef_mtd', {})
    def _pct1(v):
        if v is None: return None
        return f"{v*100:.1f}".replace('.', ',')
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text.strip()
        l = _left(sh)
        if l < 6 or ('Hedef' not in t and 'Denim' not in t): continue
        if _left(sh) > 11: continue
        paras = sh.text_frame.paragraphs
        for pi, para in enumerate(paras):
            runs = para.runs
            if not runs: continue
            all_text = ''.join(r.text for r in runs)
            matched_key = None
            for label, furkan_key in cat_map.items():
                if label.lower().split()[0] in all_text.lower():
                    matched_key = furkan_key; break
            hedef_runs = [r for r in runs if 'Hedef' in r.text]
            grc_runs   = [r for r in runs if 'Gerçekleşen' in r.text]
            if matched_key:
                g = _pct1(grc.get(matched_key))
                h = _pct1(hdf.get(matched_key))
                for r in hedef_runs:
                    prefix = r.text.split(':')[0] + ': '
                    r.text = prefix + (h if h else '?')
                for r in grc_runs:
                    prefix = r.text.split(':')[0] + ': '
                    r.text = prefix + (g if g else '?')
        # Tees özel
        for pi, para in enumerate(paras):
            runs = para.runs
            if not runs: continue
            full = ''.join(r.text for r in runs)
            if full.strip().lower() == 'tees' and pi+1 < len(paras):
                next_runs = paras[pi+1].runs
                tees_furkan = None
                for label, fk in cat_map.items():
                    if 'tees' in label.lower(): tees_furkan = fk; break
                if tees_furkan:
                    g = _pct1(grc.get(tees_furkan))
                    h = _pct1(hdf.get(tees_furkan))
                    for r in next_runs:
                        if 'Hedef' in r.text:
                            r.text = r.text.split(':')[0].rstrip()+': '+(h if h else '?')
                        elif 'Gerçekleşen' in r.text:
                            r.text = r.text.split(':')[0].rstrip()+': '+(g if g else '?')
        break


# ── Slayt 1: Kapak ───────────────────────────────────────────────────────────
def update_slide1(slide, rapor, config):
    """Kapak sayfasındaki tarih metnini günceller."""
    ay     = rapor.get('ay_etiket', config.get('month_label', ''))
    dr     = config.get('date_range', '')
    old_ay = config.get('old_ay_etiket', '')
    old_dr = config.get('old_date_range', '')
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if old_ay and old_ay in run.text:
                    run.text = run.text.replace(old_ay, ay)
                if old_dr and old_dr in run.text:
                    run.text = run.text.replace(old_dr, dr)
                # "Mayıs" kelimesini "Haziran" ile değiştir (her ihtimale karşı)
                if old_ay:
                    old_month = old_ay.split("'")[0]  # "Mayıs'26" → "Mayıs"
                    new_month = ay.split("'")[0] if ay else ''
                    if old_month and new_month and old_month in run.text:
                        run.text = run.text.replace(old_month, new_month)
