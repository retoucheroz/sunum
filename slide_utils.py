"""pptx shape manipülasyon yardımcıları."""
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

NAVY  = RGBColor(0x1F, 0x40, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x14, 0x14, 0x14)
LIGHT = RGBColor(0xDE, 0xE9, 0xF5)
GREEN = RGBColor(0x00, 0x70, 0x00)
RED   = RGBColor(0xC0, 0x00, 0x00)

NS_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

def _remove(sh):
    sh._element.getparent().remove(sh._element)

def _clear_all_runs(para):
    """Paragraftaki TÜM run'ları siler."""
    for r_el in list(para._p.findall(f'{NS_A}r')):
        para._p.remove(r_el)
    # \x0b (line break) elementlerini de temizle
    for br in list(para._p.findall(f'{NS_A}br')):
        para._p.remove(br)

def set_para_text(shape, para_idx, text, bold=None, color=None, size_pt=None):
    """Paragrafı tamamen temizleyip tek run yazar. Format ilk run'dan miras alınır."""
    tf = shape.text_frame
    if para_idx >= len(tf.paragraphs): return
    p = tf.paragraphs[para_idx]
    # Mevcut font bilgisini kaydet
    old_bold = old_color = old_size = None
    if p.runs:
        r0 = p.runs[0]
        old_bold  = r0.font.bold
        old_size  = r0.font.size
        try: old_color = r0.font.color.rgb
        except: pass
    _clear_all_runs(p)
    if not text: return
    run = p.add_run()
    run.text = text
    run.font.bold  = bold  if bold  is not None else old_bold
    run.font.size  = Pt(size_pt) if size_pt else old_size
    if color: run.font.color.rgb = color
    elif old_color:
        try: run.font.color.rgb = old_color
        except: pass

def set_bullet_para(shape, para_idx, label, value, label_bold=False,
                    value_bold=True, value_color=None):
    tf = shape.text_frame
    if para_idx >= len(tf.paragraphs): return
    p = tf.paragraphs[para_idx]
    _clear_all_runs(p)
    r1 = p.add_run(); r1.text = label; r1.font.bold = label_bold
    r2 = p.add_run(); r2.text = value; r2.font.bold = value_bold
    if value_color: r2.font.color.rgb = value_color

def find_shapes(slide, min_w_in=None, left_gt_in=None, left_lt_in=None,
                text_contains=None, shape_type=None):
    result = []
    for sh in slide.shapes:
        if shape_type and int(sh.shape_type) != shape_type: continue
        w = Emu(sh.width).inches if sh.width else 0
        l = Emu(sh.left).inches if sh.left else 0
        if min_w_in and w < min_w_in: continue
        if left_gt_in and l <= left_gt_in: continue
        if left_lt_in and l >= left_lt_in: continue
        if text_contains and sh.has_text_frame:
            if text_contains not in sh.text_frame.text: continue
        result.append(sh)
    return result

def add_table(slide, left_in, top_in, width_in, rows_data, col_widths_in,
              row_height_pt=18, header_fill=None, body_fill=None,
              header_font_color=None, body_font_color=None,
              bold_rows=None, alt_rows=None, font_size=10):
    from pptx.util import Inches
    hf  = header_fill  or RGBColor(0x1F, 0x40, 0x63)
    bf  = body_fill    or RGBColor(0xFF, 0xFF, 0xFF)
    hfc = header_font_color or RGBColor(0xFF, 0xFF, 0xFF)
    bfc = body_font_color   or RGBColor(0x14, 0x14, 0x14)
    alt_fill = RGBColor(0xDE, 0xE9, 0xF5)
    nrows = len(rows_data); ncols = len(col_widths_in)
    tbl_w = int(sum(col_widths_in) * 914400)
    tbl = slide.shapes.add_table(nrows, ncols,
                                  Emu(int(left_in*914400)), Emu(int(top_in*914400)),
                                  Emu(tbl_w),
                                  Emu(int(row_height_pt * 12700 * nrows))).table
    for ci, cw in enumerate(col_widths_in):
        tbl.columns[ci].width = Emu(int(cw*914400))
    for ri in range(nrows):
        tbl.rows[ri].height = Emu(int(row_height_pt * 12700))
    bold_rows = bold_rows or set()
    alt_rows  = alt_rows  or set()
    for ri, row_vals in enumerate(rows_data):
        is_header = ri == 0
        is_bold   = ri in bold_rows or is_header
        is_alt    = ri in alt_rows
        fill_clr  = hf if is_header else (alt_fill if is_alt else bf)
        font_clr  = hfc if is_header else bfc
        for ci, val in enumerate(row_vals[:ncols]):
            cell = tbl.cell(ri, ci)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgb.set('val', f'{fill_clr.red:02X}{fill_clr.green:02X}{fill_clr.blue:02X}')
            tf = cell.text_frame
            p = tf.paragraphs[0]
            try:
                float(str(val).replace('%','').replace(',','.'))
                p.alignment = PP_ALIGN.RIGHT
            except ValueError:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val) if val is not None else ""
            run.font.size = Pt(font_size)
            run.font.bold = is_bold
            run.font.color.rgb = font_clr
    return tbl
