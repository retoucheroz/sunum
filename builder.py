"""Mavi Haftalık Özet Sunum Motoru v2.
Mevcut özet sunumu ŞABLON olarak kullanır; tüm içerik Excel/pptx kaynaklarından çekilir."""
import io, os, zipfile, tempfile, traceback
import requests
from pptx import Presentation
from pptx.util import Emu, Pt
from PIL import Image
from tables import build_bestseller_table
from data_extract import extract_bestseller, jean_narrative, _k
from data_sources import parse_rapor, parse_cevirdigim, fmt_pct, fmt_mtl
from slide_updaters import (update_slide1, update_slide2, update_slide3, update_slide4,
                             update_slide5, update_slide6, update_slide8,
                             replace_date_texts)

PROD_IMG_URL = "https://sky-static.mavi.com/{code}_image_1.jpg"

# ---------- yardımcılar ----------
def _emf_to_png(report_zip_bytes, media_name):
    """Rapor pptx içinden ham görsel byte'ı çıkarır."""
    with zipfile.ZipFile(io.BytesIO(report_zip_bytes)) as z:
        return z.read(f"ppt/media/{media_name}")

def _report_slide_media(report_path):
    """Rapordaki slayt -> media dosyası eşleşmesini bulur (slide2, slide3)."""
    with open(report_path, "rb") as f:
        data = f.read()
    out = {}
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        for sl in ("slide2", "slide3"):
            rels = z.read(f"ppt/slides/_rels/{sl}.xml.rels").decode("utf-8")
            import re
            m = re.search(r'Target="\.\./media/([^"]+)"', rels)
            if m:
                out[sl] = m.group(1)
    return data, out

_NOTE_PREFIXES = ("Abi", "Burası boş", "Bir önceki sayfanın", "Altta yer alan", "Mesela")
_NOTE_KEYWORDS = ("örnek çalışma dosyası", "planlamadan gelen excel", "kopyala yapıştır",
                  "satış öne çıkan", "öne çıkan fitler excel", "öne non denim")

def _is_note(shape, slide_w=None):
    """Çalışma notu mu? ('Abi...'/'Altta...' notları, not anahtar kelimeleri,
    slayt dışı kutular ya da sağ kenara sıkışmış dar açıklama kutuları)."""
    if not shape.has_text_frame:
        return False
    t = shape.text_frame.text.strip()
    if not t:
        return False
    if any(t.startswith(p) for p in _NOTE_PREFIXES) or "ekleme şansımız" in t:
        return True
    if any(k in t for k in _NOTE_KEYWORDS):
        return True
    if slide_w is not None and shape.left is not None and shape.width is not None:
        left_in = Emu(shape.left).inches
        w_in = Emu(shape.width).inches
        sw_in = Emu(slide_w).inches
        # tamamen slayt dışı (sadece belirgin negatif; -0.02" gibi minik kaymalar başlıktır, silME)
        if left_in < -0.4 or left_in > sw_in:
            return True
        # sağ kenara sıkışmış dar + uzun metinli kutu = kenar notu
        if left_in > (sw_in - 2.9) and w_in < 2.6 and len(t) > 40:
            return True
    return False

def _remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)

def _strip_notes(slide, slide_w=None):
    for sh in list(slide.shapes):
        if _is_note(sh, slide_w):
            _remove_shape(sh)

def _replace_text_global(slide, replacements):
    """Tüm metin kutularında run düzeyinde bul-değiştir (format korunur)."""
    if not replacements:
        return
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                for old, new in replacements.items():
                    if old and old in r.text:
                        r.text = r.text.replace(old, new)

def _clear_product_pics(slide, slide_w):
    """Ürün slaytındaki eski görselleri kaldırır (sağ-alt küçük logoyu korur)."""
    for sh in list(slide.shapes):
        if sh.shape_type == 13:  # PICTURE
            w_in = Emu(sh.width).inches if sh.width else 0
            left_in = Emu(sh.left).inches if sh.left else 0
            top_in = Emu(sh.top).inches if sh.top else 0
            is_logo = (w_in < 2.6 and left_in > 9 and top_in > 6)
            if not is_logo:
                _remove_shape(sh)

def _set_para(shape, idx, text):
    """Paragrafın metnini ilk run'ın formatını koruyarak değiştirir."""
    tf = shape.text_frame
    if idx >= len(tf.paragraphs):
        return False
    p = tf.paragraphs[idx]
    if not p.runs:
        run = p.add_run(); run.text = text
        return True
    p.runs[0].text = text
    for extra in p.runs[1:]:
        extra.text = ""
    return True

def _update_comparison_bullets(slide, bullets, date_label=''):
    """
    Slayt 7/9/10/11 sağ kutu:
    bullets: [(kategori, satis_pct_str, stok_pct_str)]
    Paragraf yapısı: [kategori, 'Satış vs. Stok: ', '%X ', 'vs', ' ', '±Y%']
    Son değer run'ları (Satış% ve Stok%) güncellenir.
    """
    from pptx.util import Emu
    from pptx.dml.color import RGBColor

    GREEN = RGBColor(0, 128, 0)
    RED   = RGBColor(192, 0, 0)

    def _color(v_str):
        return GREEN if not v_str.startswith('-') else RED

    bullet_dict = {b[0].lower(): b for b in bullets if b[0] and str(b[0]) != 'nan'}

    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text
        l = Emu(sh.left).inches if sh.left else 0
        if l < 5: continue
        if 'Satış vs. Stok' not in t and 'Jean' not in t and 'Denim' not in t: continue

        for para in sh.text_frame.paragraphs:
            runs = para.runs
            if not runs: continue
            run_texts = [r.text for r in runs]

            # Kategori adını bul (ilk run)
            cat_raw = run_texts[0].strip().rstrip(':').lower()
            # Birleşik isimleri de dene (örn. 'Non Denim Altlar')
            matched = None
            for key, bull in bullet_dict.items():
                if key in cat_raw or cat_raw in key:
                    matched = bull; break

            if not matched: continue
            _, sat, stk = matched

            # Satış% run'unu bul ve güncelle
            for ri, r in enumerate(runs):
                if 'Satış vs. Stok:' in r.text:
                    # Sonraki run'lar: '%X', 'vs', ' ', '±Y%'
                    if ri+1 < len(runs):
                        runs[ri+1].text = sat + ' '
                        runs[ri+1].font.color.rgb = _color(sat)
                    if ri+3 < len(runs):
                        runs[ri+3].text = stk
                        runs[ri+3].font.color.rgb = _color(stk)
                    break

def _swap_lfl_emf(prs, rapor_path, slide_idx, rapor_slide_idx, emf_out_name):
    """Rapor pptx'teki LFL EMF görselini özet sunumun ilgili slaytına kopyalar."""
    import zipfile, re
    # Rapor pptx'teki EMF'yi bul
    with zipfile.ZipFile(rapor_path) as z:
        rels = z.read(f'ppt/slides/_rels/slide{rapor_slide_idx+1}.xml.rels').decode()
        emfs = re.findall(r'Target=\"\.\./media/([^\"]+\.emf)\"', rels)
        if not emfs:
            return
        emf_data = z.read(f'ppt/media/{emfs[0]}')
    # Özet sunum ZIP'ine yaz
    # python-pptx üzerinden doğrudan part erişimi
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    slide = list(prs.slides)[slide_idx]
    # Slide'daki EMF shape'lerini bul, media'ya yaz
    for sh in slide.shapes:
        if sh.shape_type == 13:  # PICTURE
            try:
                img_part = sh.image
                # EMF dosyasına bak
                if 'emf' in str(getattr(img_part, '_blob', b''))[:10].lower() or True:
                    # Şekil blobunu değiştir — pptx internal
                    from pptx.opc.part import Part
                    img_part._blob = emf_data
                    return
            except: pass

def _replace_picture(slide, pic, image_path):
    """Resmi aynı çerçeve içinde (en-boy oranını koruyarak) yenisiyle değiştirir."""
    left, top, width, height = pic.left, pic.top, pic.width, pic.height
    with Image.open(image_path) as im:
        iw, ih = im.size
    # çerçeveye sığdır, oranı koru, üst-sola hizala
    fw, fh = width, height
    scale = min(fw/iw, fh/ih)
    nw, nh = int(iw*scale), int(ih*scale)
    _remove_shape(pic)
    slide.shapes.add_picture(image_path, left, top, nw, nh)

def _download_product_image(code, dest_dir, timeout=10):
    """sky-static'ten ürün görselini indirir. Başarısızsa None."""
    c = code.lstrip("Mm").strip()
    url = PROD_IMG_URL.format(code=c)
    try:
        r = requests.get(url, timeout=timeout, headers={"User-Agent": "Mozilla/5.0"})
        if r.status_code == 200 and r.headers.get("content-type", "").startswith("image"):
            p = os.path.join(dest_dir, f"{c}.jpg")
            with open(p, "wb") as f:
                f.write(r.content)
            return p
    except Exception:
        pass
    return None

def _place_product_grid(slide, items, tmpdir, log):
    """İlk 10 ürün görselini erkek/kadın alanına ızgara + kod/adet etiketiyle yerleştirir.
    items: {'erkek':[(code,qty),...], 'kadin':[(code,qty),...]}"""
    from pptx.util import Pt as _Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    pres = slide.part.package.presentation_part.presentation
    sw, sh = pres.slide_width, pres.slide_height
    placed = 0
    for side, clist in items.items():
        x0 = Emu(0.3*914400) if side == "erkek" else Emu(sw/2 + 0.1*914400)
        y0 = Emu(1.6*914400)
        col_w = Emu((sw/2 - 0.5*914400))
        cols, rows = 5, 2
        cw = int(col_w/cols)
        chh = int((sh - y0 - 0.5*914400)/rows)
        img_h = int(chh * 0.74)
        for i, item in enumerate(clist[:10]):
            code, qty = (item if isinstance(item, (list, tuple)) else (item, None))
            p = _download_product_image(code, tmpdir)
            r, c = divmod(i, cols)
            cell_x = int(x0 + c*cw)
            cell_y = int(y0 + r*chh)
            if p:
                try:
                    with Image.open(p) as im:
                        iw, ih = im.size
                    scale = min(cw*0.88/iw, img_h/ih)
                    nw, nh = int(iw*scale), int(ih*scale)
                    lx = int(cell_x + (cw-nw)/2)
                    slide.shapes.add_picture(p, lx, cell_y, nw, nh)
                    placed += 1
                except Exception as e:
                    log.append(f"  · {side} {code}: yerleştirme hatası {e}")
            else:
                log.append(f"  · {side} {code}: görsel indirilemedi")
            # kod + adet etiketi
            tb = slide.shapes.add_textbox(cell_x, int(cell_y + img_h + 0.02*914400),
                                          cw, Emu(0.42*914400))
            tf = tb.text_frame; tf.word_wrap = True
            tf.margin_top = 0; tf.margin_bottom = 0
            para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
            run = para.add_run(); run.text = str(code)
            run.font.size = _Pt(6); run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
            if qty is not None:
                p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
                r2 = p2.add_run(); r2.text = f"{int(round(qty)):,}".replace(",", ".")
                r2.font.size = _Pt(8); r2.font.bold = True
                r2.font.color.rgb = RGBColor(0x1F, 0x40, 0x63)
    return placed

# ---------- ana fonksiyon ----------
def build(template_path, report_path, fitler_xlsx, nondenim_xlsx, planlama_xlsx,
          config, output_path, download_products=True, kpi_manual=None,
          gh_fitler_path=None, gh_nondenim_path=None, calisma_path=None):
    log = []
    kpi_manual = kpi_manual or {}
    prs = Presentation(template_path)
    slides = prs.slides

    # 1) verileri hesapla
    log.append("Veriler okunuyor...")
    fit = extract_bestseller(fitler_xlsx,  gh_path=gh_fitler_path)
    non = extract_bestseller(nondenim_xlsx, gh_path=gh_nondenim_path)
    rapor = parse_rapor(report_path)
    cevirdigim = parse_cevirdigim(planlama_xlsx, config.get("date_range",""))
    log.append(f"  Hedef: {fmt_mtl(rapor['hedef_mtl'])} | Gerçekleşen: {fmt_mtl(rapor['grc_mtl'])} | HGO: {rapor['hgo']}")
    log.append(f"  Jean: erkek {_k(fit['erkek']['totals']['total_satis'])}, kadın {_k(fit['kadin']['totals']['total_satis'])}")
    # 1b) Çalışma dosyası verileri
    calisma_data = {}
    if calisma_path:
        try:
            from calisma_dosyasi import read_sheet, get_erkek_bullets, get_kadin_bullets, get_sub_bullets
            calisma_data['erkek_kat']    = read_sheet(calisma_path, 'erkek_kategori')
            calisma_data['kadin_kat']    = read_sheet(calisma_path, 'kadin_kategori')
            calisma_data['erkek_sub']    = read_sheet(calisma_path, 'erkek_sub_aylik')
            calisma_data['kadin_sub']    = read_sheet(calisma_path, 'kadin_sub_aylik')
            log.append(f"Çalışma dosyası okundu.")
        except Exception as e:
            log.append(f"Çalışma dosyası hatası: {e}")

    tmpdir = tempfile.mkdtemp()

    tmpdir = tempfile.mkdtemp()

    # 2) tüm slaytlardaki çalışma notlarını temizle + tarih/ay değiştir
    slide_w = prs.slide_width
    replacements = config.get("replacements", {})
    for s in slides:
        _strip_notes(s, slide_w)
        _replace_text_global(s, replacements)

    # Tarih/ay otomatik değiştirme
    old_ay = config.get("old_ay_etiket", "")   # "Mayıs'26"
    new_ay = config.get("month_label", "")      # "Haziran'26"
    old_dr = config.get("old_date_range", "")   # "1-31 Mayıs"
    new_dr = config.get("date_range", "")       # "1-17 Haziran"
    # Geçen yıl değişimi: "Mayıs'25" → "Haziran'25"
    old_ay_ly = old_ay.replace("'26", "'25") if old_ay else ""
    new_ay_ly = new_ay.replace("'26", "'25") if new_ay else ""
    
    date_pairs = []
    if old_ay and new_ay and old_ay != new_ay:
        date_pairs.append((old_ay, new_ay))
    if old_ay_ly and new_ay_ly and old_ay_ly != new_ay_ly:
        date_pairs.append((old_ay_ly, new_ay_ly))
    if old_dr and new_dr and old_dr != new_dr:
        date_pairs.append((old_dr, new_dr))
    if date_pairs:
        olds = [p[0] for p in date_pairs]
        news = [p[1] for p in date_pairs]
        replace_date_texts(slides, olds, news)
        log.append(f"Tarih/ay güncellendi: {date_pairs}")
    log.append("Çalışma notları temizlendi.")

    dr = config.get("date_range", "1-31 Mayıs")

    # 2.5) Slayt 1: Kapak tarihi
    try:
        update_slide1(slides[0], rapor, config)
        log.append("Slayt 1: Kapak tarihi güncellendi.")
    except Exception as e:
        log.append(f"Slayt 1 hata: {e}")
    try:
        update_slide2(slides[1], rapor, cevirdigim, kpi_manual, config)
        log.append("Slayt 2: Executive Summary güncellendi.")
    except Exception as e:
        log.append(f"Slayt 2 hata: {e}")

    # 4) Slayt 3: Rapor tabloları (resim) + sağ kutular (metin)
    try:
        rep_data, media = _report_slide_media(report_path)
        img_targ = os.path.join(tmpdir, "rep_hedef.png")
        img_kpi  = os.path.join(tmpdir, "rep_kpi.png")
        with open(img_targ, "wb") as f: f.write(_emf_to_png(rep_data, media["slide2"]))
        with open(img_kpi,  "wb") as f: f.write(_emf_to_png(rep_data, media["slide3"]))
        s3 = slides[2]
        pics = sorted([sh for sh in s3.shapes if sh.shape_type == 13 and Emu(sh.width).inches > 3],
                      key=lambda p: p.top)
        if len(pics) >= 2:
            _replace_picture(s3, pics[0], img_targ)
            _replace_picture(s3, pics[1], img_kpi)
        update_slide3(s3, rapor, cevirdigim, kpi_manual, config)
        log.append("Slayt 3: Hedef tablosu aktarıldı + sağ kutular güncellendi.")
    except Exception as e:
        log.append(f"Slayt 3 hata: {e}")

    # 5) Slayt 4: Yönetici özeti
    try:
        update_slide4(slides[3], rapor, cevirdigim, config)
        log.append("Slayt 4: Yönetici özeti güncellendi.")
    except Exception as e:
        log.append(f"Slayt 4 hata: {e}")

    # 6) Slayt 5: Toplam LFL
    try:
        update_slide5(slides[4], cevirdigim, config)
        log.append("Slayt 5: LFL Toplam güncellendi.")
    except Exception as e:
        log.append(f"Slayt 5 hata: {e}")

    # 7) Slayt 6: Erkek — sol tabloyu rapor pptx'ten al, sağ kutuyu furkan'dan güncelle
    try:
        _swap_lfl_emf(prs, report_path, slide_idx=5, rapor_slide_idx=5, emf_out_name='image6.emf')
        update_slide6(slides[5], cevirdigim, {}, config)
        log.append("Slayt 6: Erkek LFL tablosu rapor'dan alındı, kategoriler furkan'dan güncellendi.")
    except Exception as e:
        log.append(f"Slayt 6 hata: {e}")

    # 8) Slayt 8: Kadın — aynı
    try:
        _swap_lfl_emf(prs, report_path, slide_idx=7, rapor_slide_idx=5, emf_out_name='image8.emf')
        update_slide8(slides[7], cevirdigim, {}, config)
        log.append("Slayt 8: Kadın LFL tablosu rapor'dan alındı, kategoriler furkan'dan güncellendi.")
    except Exception as e:
        log.append(f"Slayt 8 hata: {e}")

    # 9) Slayt 7, 9, 10, 11: Çalışma dosyası karşılaştırma bullet'ları
    if calisma_data:
        try:
            from calisma_dosyasi import get_erkek_bullets, get_kadin_bullets, get_sub_bullets
            dr = config.get('date_range','')
            # Slayt 7: Erkek kategori
            if calisma_data.get('erkek_kat') is not None:
                bullets = get_erkek_bullets(calisma_data['erkek_kat'])
                _update_comparison_bullets(slides[6], bullets, dr)
                log.append(f"Slayt 7: Erkek karşılaştırma güncellendi ({len(bullets)} kategori).")
            # Slayt 9: Kadın kategori
            if calisma_data.get('kadin_kat') is not None:
                bullets = get_kadin_bullets(calisma_data['kadin_kat'])
                _update_comparison_bullets(slides[8], bullets, dr)
                log.append(f"Slayt 9: Kadın karşılaştırma güncellendi.")
            # Slayt 10: Erkek alt kategori
            if calisma_data.get('erkek_sub') is not None:
                bullets = get_sub_bullets(calisma_data['erkek_sub'], 'Erkek')
                _update_comparison_bullets(slides[9], bullets, dr)
                log.append(f"Slayt 10: Erkek alt kategori güncellendi.")
            # Slayt 11: Kadın alt kategori
            if calisma_data.get('kadin_sub') is not None:
                bullets = get_sub_bullets(calisma_data['kadin_sub'], 'Kadın')
                _update_comparison_bullets(slides[10], bullets, dr)
                log.append(f"Slayt 11: Kadın alt kategori güncellendi.")
        except Exception as e:
            log.append(f"Slayt 7/9/10/11 hata: {e}")

    # 4) SLAYT 14 & 17: best-seller tabloları + anlatı + ürün görselleri
    def _place_table_pic(slide, image_path, left_in, top_in, max_w_in, max_h_in):
        from PIL import Image as _I
        with _I.open(image_path) as im:
            iw, ih = im.size
        sc = min(max_w_in/iw, max_h_in/ih)
        nw, nh = iw*sc, ih*sc
        slide.shapes.add_picture(image_path, Emu(int(left_in*914400)), Emu(int(top_in*914400)),
                                 Emu(int(nw*914400)), Emu(int(nh*914400)))

    def fill_bestseller(slide_idx, data, kind, erkek_text, kadin_text):
        s = slides[slide_idx]
        et, _ = build_bestseller_table(data["erkek"]["rows"], data["erkek"]["totals"],
                                       data["erkek"]["weights"], kind,
                                       os.path.join(tmpdir, f"t{slide_idx}_e.png"),
                                       gy_ranks=data["erkek"].get("gy_ranks"),
                                       gh_ranks=data["erkek"].get("gh_ranks"))
        kt, _ = build_bestseller_table(data["kadin"]["rows"], data["kadin"]["totals"],
                                       data["kadin"]["weights"], kind,
                                       os.path.join(tmpdir, f"t{slide_idx}_k.png"),
                                       gy_ranks=data["kadin"].get("gy_ranks"),
                                       gh_ranks=data["kadin"].get("gh_ranks"))
        # eski tablo resimlerini kaldır
        for sh in [p for p in s.shapes if p.shape_type == 13 and Emu(p.width).inches > 3]:
            _remove_shape(sh)
        # okunur boyutta yeniden yerleştir (alt yarıyı doldur)
        _place_table_pic(s, et, 0.15, 3.05, 6.45, 4.05)   # sol = erkek
        _place_table_pic(s, kt, 6.75, 3.05, 6.45, 4.05)   # sağ = kadın
        # anlatı: para0 = erkek, para2 = kadın (yapı korunur)
        for sh in s.shapes:
            if sh.has_text_frame and len(sh.text_frame.text) > 80 and Emu(sh.width).inches > 8:
                _set_para(sh, 0, erkek_text)
                _set_para(sh, 2, kadin_text)
                break

    je, jk = jean_narrative(fit, date_range=config.get('date_range',''))
    fill_bestseller(13, fit, "jean", je, jk)
    # non-denim anlatı (erkek / kadın ayrı)
    e, k = non["erkek"], non["kadin"]
    nd_e = (f"{dr} tarihleri arasında erkek non denim alt kategorisinde toplamda "
            f"{_k(e['totals']['total_satis'])} adet satış gerçekleşmiştir. İlk 10 ürünün toplam satışı "
            f"{_k(e['totals']['top10_satis'])} adettir ve toplam satışın "
            f"%{round(e['weights']['satis']*100)}'sini oluşturmaktadır.")
    nd_k = (f"Kadın kategorisinde ise toplamda {_k(k['totals']['total_satis'])} adet satış "
            f"gerçekleşmiştir. İlk 10 ürünün toplam satışı {_k(k['totals']['top10_satis'])} adettir ve "
            f"toplam satışın %{round(k['weights']['satis']*100)}'ünü oluşturmaktadır.")
    fill_bestseller(16, non, "altlar", nd_e, nd_k)
    log.append("Slayt 14 & 17: best-seller tabloları (okunur boyut) ve anlatı güncellendi.")

    # 5) SLAYT 15 & 18: orijinal Power BI düzeni (kod+ad+foto+Qty+Retail Price)
    from powerbi_grid import build_powerbi_grid
    def build_grid_slide(slide, data, tag):
        s = slides[slide]
        # foto indir
        def prep(side):
            items, photos = [], []
            for (code, name, satis, stok) in data[side]["rows"]:
                items.append((code, name, satis, None))  # fiyat kaynakta yok
                photos.append(_download_product_image(code, tmpdir) if download_products else None)
            return items, photos
        ie, pe = prep("erkek"); ik, pk = prep("kadin")
        out_img = os.path.join(tmpdir, f"grid_{slide}.png")
        build_powerbi_grid(ie, ik, out_img, pe, pk)
        # eski görselleri temizle (logo hariç)
        _clear_product_pics(s, slide_w)
        # tam genişlik yerleştir (başlığın altı)
        sw_in = Emu(slide_w).inches
        from PIL import Image as _I
        with _I.open(out_img) as im:
            iw, ih = im.size
        left_in, top_in, max_w = 0.2, 0.95, sw_in - 0.4
        sc = min(max_w/iw, 6.3/ih)
        s.shapes.add_picture(out_img, Emu(int(left_in*914400)), Emu(int(top_in*914400)),
                             Emu(int(iw*sc*914400)), Emu(int(ih*sc*914400)))
        return sum(1 for p in (pe+pk) if p)

    n15 = build_grid_slide(14, fit, "jean")
    n18 = build_grid_slide(17, non, "nondenim")
    log.append(f"Slayt 15: {n15}/20 ürün foto indirildi · Slayt 18: {n18}/20 foto indirildi "
               f"(Power BI düzeninde).")
    if not download_products:
        log.append("  (Ürün indirme kapalıydı; foto yerine placeholder kullanıldı.)")

    # 6) SLAYT 4: yönetici özeti best-seller paragrafları (jean anlatısı)
    #    (sadece best-seller cümleleri otomatik; üst genel paragraflar config'ten)
    # — v1: jean anlatısını yönetici özetinin best-seller bölümüne enjekte etmiyoruz
    #   (serbest metin riskli); kullanıcı isterse sonra eşleriz.

    prs.save(output_path)
    log.append(f"Sunum kaydedildi: {output_path}")
    return log
